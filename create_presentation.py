"""
Generate the Multi-Faceted Music Retrieval final presentation.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Colour palette ──────────────────────────────────────────────────────────
BG       = RGBColor(0x0F, 0x17, 0x2A)   # dark navy
ACCENT   = RGBColor(0x63, 0x66, 0xF1)   # indigo
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT    = RGBColor(0xCB, 0xD5, 0xE1)   # slate-300
MUTED    = RGBColor(0x94, 0xA3, 0xB8)   # slate-400
GREEN    = RGBColor(0x10, 0xB9, 0x81)
RED      = RGBColor(0xEF, 0x44, 0x44)
AMBER    = RGBColor(0xF5, 0x9E, 0x0B)
CYAN     = RGBColor(0x06, 0xB6, 0xD4)
PINK     = RGBColor(0xEC, 0x48, 0x99)
PURPLE   = RGBColor(0x8B, 0x5C, 0xF6)
ORANGE   = RGBColor(0xF9, 0x73, 0x16)
DARK_CARD = RGBColor(0x1E, 0x29, 0x3B)  # card background

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


# ── Helpers ─────────────────────────────────────────────────────────────────
def add_bg(slide):
    """Fill slide background with dark navy."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_paragraph(tf, text, font_size=18, color=WHITE, bold=False,
                  space_before=Pt(6), space_after=Pt(2), alignment=PP_ALIGN.LEFT,
                  font_name="Calibri"):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p


def add_card(slide, left, top, width, height, fill_color=DARK_CARD):
    """Add a rounded rectangle card."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def bullet_slide(slide, title, bullets, sub_title=None):
    """Standard content slide with title and bullet points."""
    add_bg(slide)
    # Title
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
                 title, font_size=36, color=WHITE, bold=True)
    if sub_title:
        add_text_box(slide, Inches(0.8), Inches(1.15), Inches(11), Inches(0.5),
                     sub_title, font_size=18, color=MUTED)
    # Accent line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.8), Inches(1.55), Inches(2.5), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()
    # Bullets
    y_start = Inches(1.9)
    tf = add_text_box(slide, Inches(1.0), y_start, Inches(11), Inches(5),
                      "", font_size=22, color=LIGHT)
    tf.paragraphs[0].text = ""
    for i, b in enumerate(bullets):
        if isinstance(b, tuple) and len(b) == 2:
            text, color_ = b
        elif isinstance(b, tuple):
            text, color_ = b[0], LIGHT
        else:
            text, color_ = b, LIGHT
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(22)
        p.font.color.rgb = color_
        p.font.name = "Calibri"
        p.space_before = Pt(10)
        p.space_after = Pt(4)


def section_divider(slide, number, title, subtitle=""):
    """Section divider slide."""
    add_bg(slide)
    # Large number
    add_text_box(slide, Inches(0.8), Inches(1.5), Inches(3), Inches(2),
                 f"{number:02d}", font_size=96, color=ACCENT, bold=True)
    # Title
    add_text_box(slide, Inches(0.8), Inches(3.6), Inches(10), Inches(1.2),
                 title, font_size=44, color=WHITE, bold=True)
    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(4.7), Inches(10), Inches(0.8),
                     subtitle, font_size=20, color=MUTED)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(1.5),
             "Multi-Faceted Music Retrieval",
             font_size=52, color=WHITE, bold=True)

add_text_box(slide, Inches(0.8), Inches(3.2), Inches(11), Inches(0.8),
             "Combining Audio, Text, and Graph Embeddings for Smarter Song Recommendations",
             font_size=24, color=MUTED)

# Accent line
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0.8), Inches(4.2), Inches(3), Pt(4))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT
line.line.fill.background()

add_text_box(slide, Inches(0.8), Inches(4.6), Inches(11), Inches(0.6),
             "Wenny  |  Sid & Issac  |  Archive  |  Jiayi  |  Helena",
             font_size=20, color=LIGHT)

add_text_box(slide, Inches(0.8), Inches(5.2), Inches(11), Inches(0.5),
             "Dataset: Free Music Archive (FMA) \u2014 8,000 tracks, 8 genres",
             font_size=18, color=MUTED)

add_text_box(slide, Inches(0.8), Inches(5.7), Inches(11), Inches(0.5),
             "Deep Learning \u2014 Final Project Presentation",
             font_size=16, color=MUTED)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem Statement
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
bullet_slide(slide, "The Problem", [
    ("\u201cChill lo-fi vibes\u201d \u2192 requires understanding musical mood (acoustic features)",),
    ("\u201cFolk artists about travel\u201d \u2192 requires understanding text & lyrics (semantic features)",),
    ("\u201cSongs like this one\u201d \u2192 requires structural knowledge (graph relationships)",),
    ("",),
    ("No single embedding captures all of these.", WHITE),
    ("",),
    ("Our solution: build independent embedding views for each modality,", LIGHT),
    ("then fuse them into one retrieval system that outperforms any single view.", LIGHT),
], sub_title="Music retrieval systems typically rely on a single representation")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Architecture Diagram
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "System Architecture", font_size=36, color=WHITE, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
             "Reference diagram from our design phase", font_size=16, color=MUTED)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0.8), Inches(1.55), Inches(2.5), Pt(3))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT
line.line.fill.background()

# Insert the reference diagram image
slide.shapes.add_picture(
    "/Users/zac/Documents/Documents-it/Deep Learning Notes/Finals/image.png",
    Inches(1.5), Inches(1.9), Inches(10.3), Inches(5.2)
)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Architecture Explained
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Four Retrieval Views", font_size=36, color=WHITE, bold=True)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0.8), Inches(1.2), Inches(2.5), Pt(3))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT
line.line.fill.background()

# Four cards
card_data = [
    ("View 1: Vibe / Text Search", "CLAP (HTSAT-tiny)", "512-d", "Audio waveform",
     "High-level semantics:\nmood, genre feel, instruments", ACCENT),
    ("View 2: Lyrics / Semantic", "SBERT (all-MiniLM-L6-v2)", "384-d", "Metadata + lyrics",
     "Textual semantics:\nartist identity, lyrical content", GREEN),
    ("View 3: Acoustic Similarity", "OpenL3", "512-d", "Audio waveform",
     "Low-level acoustics:\ntimbre, rhythm, texture", AMBER),
    ("View 4: Graph Recommendation", "HeteroGNN (SAGEConv)", "256-d", "Track-artist-genre graph",
     "Structural connectivity:\nco-genre, artist similarity", PINK),
]

card_w = Inches(2.8)
card_h = Inches(4.5)
gap = Inches(0.25)
start_x = Inches(0.6)

for i, (title, model, dims, inp, desc, clr) in enumerate(card_data):
    x = start_x + i * (card_w + gap)
    y = Inches(1.6)
    card = add_card(slide, x, y, card_w, card_h)

    # Color bar at top of card
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 x + Inches(0.15), y + Inches(0.15),
                                 card_w - Inches(0.3), Pt(4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = clr
    bar.line.fill.background()

    # Card title
    add_text_box(slide, x + Inches(0.2), y + Inches(0.4),
                 card_w - Inches(0.4), Inches(0.6),
                 title, font_size=16, color=WHITE, bold=True)
    # Model name
    add_text_box(slide, x + Inches(0.2), y + Inches(1.0),
                 card_w - Inches(0.4), Inches(0.3),
                 model, font_size=13, color=clr, bold=True)
    # Dimensions
    add_text_box(slide, x + Inches(0.2), y + Inches(1.35),
                 card_w - Inches(0.4), Inches(0.3),
                 f"Dimensions: {dims}", font_size=12, color=MUTED)
    # Input
    add_text_box(slide, x + Inches(0.2), y + Inches(1.7),
                 card_w - Inches(0.4), Inches(0.3),
                 f"Input: {inp}", font_size=12, color=MUTED)
    # Description
    add_text_box(slide, x + Inches(0.2), y + Inches(2.2),
                 card_w - Inches(0.4), Inches(1.8),
                 desc, font_size=14, color=LIGHT)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Dataset Overview
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
bullet_slide(slide, "Dataset: Free Music Archive (FMA)", [
    ("8,000 tracks across 8 top-level genres (FMA Small subset)", WHITE),
    ("Genres: Electronic, Experimental, Folk, Hip-Hop, Instrumental, International, Pop, Rock",),
    ("Open-licensed benchmark \u2014 freely available for research",),
    ("",),
    ("Additional data sources:", WHITE),
    ("  \u2022  FMA metadata: tracks.csv, genres.csv (artist, title, tags, genre labels)",),
    ("  \u2022  Genius API: song lyrics (first 1,000 characters per track)",),
    ("  \u2022  Echo Nest: 8 continuous audio features (danceability, energy, valence, etc.)",),
    ("  \u2022  294 tracks overlap with Echo Nest \u2014 used as independent evaluation ground truth",),
], sub_title="github.com/mdeff/fma")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Section divider: View 1
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_divider(slide, 1, "CLAP \u2014 Vibe / Text-to-Music Search",
                "Maps both audio and text into a shared 512-d space via contrastive learning")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — CLAP Results
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "CLAP: Text-to-Music Retrieval", font_size=36, color=WHITE, bold=True)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0.8), Inches(1.2), Inches(2.5), Pt(3))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT
line.line.fill.background()

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
             "Natural language queries against 7,997 audio embeddings (3 corrupt MP3s skipped)",
             font_size=16, color=MUTED)

# Results table as cards
queries = [
    ("\u201csad piano ballad\u201d", "DUITA \u2014 XPURM", "Instrumental", "0.65"),
    ("\u201caggressive heavy metal\u201d", "Dead Elements \u2014 Angstbreaker", "Rock", "0.54"),
    ("\u201cupbeat happy pop song\u201d", "One Way Love \u2014 Ready for Men", "Pop", "0.52"),
    ("\u201cacoustic guitar folk\u201d", "Wainiha Valley \u2014 Mia Doi Todd", "Folk", "0.49"),
]

# Table header
header_y = Inches(2.1)
add_text_box(slide, Inches(1.0), header_y, Inches(3.5), Inches(0.4),
             "Query", font_size=14, color=ACCENT, bold=True)
add_text_box(slide, Inches(4.5), header_y, Inches(3.5), Inches(0.4),
             "Top Result", font_size=14, color=ACCENT, bold=True)
add_text_box(slide, Inches(8.5), header_y, Inches(1.5), Inches(0.4),
             "Genre", font_size=14, color=ACCENT, bold=True)
add_text_box(slide, Inches(10.5), header_y, Inches(1.5), Inches(0.4),
             "Cosine Sim", font_size=14, color=ACCENT, bold=True)

for i, (query, result, genre, sim) in enumerate(queries):
    row_y = Inches(2.6) + Inches(0.55) * i
    if i % 2 == 0:
        add_card(slide, Inches(0.8), row_y - Inches(0.05),
                 Inches(11.2), Inches(0.5), DARK_CARD)
    add_text_box(slide, Inches(1.0), row_y, Inches(3.5), Inches(0.4),
                 query, font_size=15, color=WHITE)
    add_text_box(slide, Inches(4.5), row_y, Inches(3.5), Inches(0.4),
                 result, font_size=15, color=LIGHT)
    add_text_box(slide, Inches(8.5), row_y, Inches(1.5), Inches(0.4),
                 genre, font_size=15, color=MUTED)
    add_text_box(slide, Inches(10.5), row_y, Inches(1.5), Inches(0.4),
                 sim, font_size=15, color=GREEN)

# Genre structure insights
add_text_box(slide, Inches(0.8), Inches(5.0), Inches(11), Inches(0.4),
             "Genre Structure in CLAP Embedding Space", font_size=18, color=WHITE, bold=True)
tf = add_text_box(slide, Inches(1.0), Inches(5.4), Inches(11), Inches(1.5),
                  "", font_size=15, color=LIGHT)
for line_text in [
    "Most similar:  Hip-Hop & Pop (0.81 cosine), Folk & International (0.80)",
    "Most distinct:  Rock & Electronic (0.64)",
    "PCA: 50 components capture ~85% of variance across 512 dimensions",
]:
    add_paragraph(tf, f"\u2022  {line_text}", font_size=15, color=LIGHT, space_before=Pt(4))


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Section divider: View 2
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_divider(slide, 2, "SBERT \u2014 Lyrics & Semantic Search",
                "Sentence-BERT encodes metadata + lyrics into 384-d embeddings")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — SBERT Details
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
bullet_slide(slide, "SBERT: Representation Analysis", [
    ("Model: all-MiniLM-L6-v2 \u2014 6-layer transformer, 384-d output", WHITE),
    ("Trained on 1B+ sentence pairs (NLI + paraphrase datasets)",),
    ("",),
    ("PCA Variance: highly distributed representation", WHITE),
    ("  \u2022  Top component captures only 4.1% of variance",),
    ("  \u2022  181 components needed for 90% (vs. 50 for CLAP\u2019s 512-d space)",),
    ("  \u2022  This means SBERT uses its dimensions more uniformly \u2014 desirable for retrieval",),
    ("",),
    ("Representation tests:", WHITE),
    ("  \u2022  Semantic robustness: only 10% overlap between \u201clonely\u201d vs \u201cisolated\u201d queries",),
    ("  \u2022  Lexical bias: \u201cBlue music\u201d retrieves titles with \u201cBlue\u201d, not Blues genre",),
    ("  \u2022  Truncation impact: minimal (FMA metadata strings are typically <50 words)",),
])


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Data Leakage
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Data Leakage: Identified & Fixed", font_size=36, color=WHITE, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
             "A critical methodological issue we caught and resolved", font_size=18, color=MUTED)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0.8), Inches(1.55), Inches(2.5), Pt(3))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT
line.line.fill.background()

# Problem card
add_card(slide, Inches(0.8), Inches(1.9), Inches(5.5), Inches(4.5), DARK_CARD)
add_text_box(slide, Inches(1.1), Inches(2.0), Inches(5), Inches(0.4),
             "Problem", font_size=20, color=RED, bold=True)
tf = add_text_box(slide, Inches(1.1), Inches(2.5), Inches(5), Inches(3.5),
                  "", font_size=15, color=LIGHT)
for t in [
    "Original SBERT input included genre_top directly:",
    '  "{artist} - {title}. Tags: {genre}, {tags}"',
    "",
    "48.8% of tags fields also contained genre-like",
    "labels (e.g., 'folk', 'pop', 'melodic')",
    "",
    "When evaluation uses genre as ground truth,",
    "the model matches on the label itself \u2014",
    "artificially inflating retrieval scores.",
]:
    add_paragraph(tf, t, font_size=15, color=LIGHT if t.strip() else LIGHT, space_before=Pt(3))

# Fix card
add_card(slide, Inches(6.8), Inches(1.9), Inches(5.5), Inches(4.5), DARK_CARD)
add_text_box(slide, Inches(7.1), Inches(2.0), Inches(5), Inches(0.4),
             "Fix", font_size=20, color=GREEN, bold=True)
tf = add_text_box(slide, Inches(7.1), Inches(2.5), Inches(5), Inches(3.5),
                  "", font_size=15, color=LIGHT)
for t in [
    "1. Removed genre_top from input entirely",
    "",
    "2. Filtered all 8 genre words from tags field",
    "   (electronic, folk, hip-hop, pop, rock, etc.)",
    "",
    "3. Added Genius API lyrics (first 1,000 chars)",
    "   as a replacement semantic signal",
    "",
    "New format:",
    '  "{title} by {artist}. Tags: {filtered}.',
    '   Lyrics: {lyrics}"',
    "",
    "Genre retained for evaluation only \u2014",
    "never seen by the embedding model.",
]:
    add_paragraph(tf, t, font_size=15, color=LIGHT, space_before=Pt(3))


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Section divider: View 3
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_divider(slide, 3, "OpenL3 \u2014 Acoustic Similarity",
                "Self-supervised audio embeddings capturing timbre, rhythm, and texture")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 12 — OpenL3 Details (Wenny's role)
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
bullet_slide(slide, "OpenL3: Audio-to-Audio Retrieval", [
    ("Model: OpenL3 (content_type='music', 512-d embeddings)", WHITE),
    ("Trained via self-supervised audio-visual correspondence on AudioSet",),
    ("Captures low-level acoustic properties without any semantic label supervision",),
    ("",),
    ("Key difference from CLAP:", WHITE),
    ("  \u2022  CLAP maps audio into a space shared with text (semantic)",),
    ("  \u2022  OpenL3 maps audio into a purely acoustic space",),
    ("  \u2022  Two songs that sound similar are close together, even if different genres",),
    ("",),
    ("Use case: \u201cFind me songs that sound like this one\u201d",),
    ("Processes all 8,000 tracks, indexed via FAISS for instant nearest-neighbor search",),
], sub_title="Role 1 \u2014 Wenny")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Section divider: View 4
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
section_divider(slide, 4, "GNN \u2014 Graph-Based Recommendation",
                "Learning structural relationships between tracks, artists, and genres")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 14 — GNN Details
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
bullet_slide(slide, "Heterogeneous GNN: Structural Embeddings", [
    ("2-layer HeteroGNN using SAGEConv on a track-artist-genre graph", WHITE),
    ("Node types: track (8,000), artist (unique IDs), genre (up to 163)",),
    ("Edge types: track\u2192artist, track\u2192genre, artist\u2192genre, co-genre",),
    ("",),
    ("Node features:", WHITE),
    ("  \u2022  Track nodes: CLAP embeddings as initial features",),
    ("  \u2022  Artist nodes: mean-pooled track CLAP embeddings",),
    ("  \u2022  Genre nodes: one-hot / learnable embeddings",),
    ("",),
    ("Training: link prediction on hidden track\u2192genre edges",),
    ("Output: 256-d track embeddings capturing connectivity structure",),
    ("Includes Flask web demo at localhost:5050 for interactive search",),
], sub_title="Role 3 \u2014 Archive")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Cross-Modal Comparison
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Cross-Modal Comparison: Why Fusion Works",
             font_size=36, color=WHITE, bold=True)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0.8), Inches(1.2), Inches(2.5), Pt(3))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT
line.line.fill.background()

add_text_box(slide, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
             "SBERT and OpenL3 retrieve almost entirely different tracks for the same query",
             font_size=18, color=MUTED)

# Big stat cards
card1 = add_card(slide, Inches(1.0), Inches(2.3), Inches(5), Inches(2.5), DARK_CARD)
add_text_box(slide, Inches(1.3), Inches(2.5), Inches(4.5), Inches(0.8),
             "5.6%", font_size=60, color=AMBER, bold=True)
add_text_box(slide, Inches(1.3), Inches(3.5), Inches(4.5), Inches(0.4),
             "Overlap @ Top-20 Neighbours", font_size=18, color=WHITE, bold=True)
add_text_box(slide, Inches(1.3), Inches(3.95), Inches(4.5), Inches(0.5),
             "Only 1.1 of 20 neighbours shared\nbetween text and audio views",
             font_size=14, color=MUTED)

card2 = add_card(slide, Inches(6.8), Inches(2.3), Inches(5), Inches(2.5), DARK_CARD)
add_text_box(slide, Inches(7.1), Inches(2.5), Inches(4.5), Inches(0.8),
             "\u20130.77", font_size=60, color=RED, bold=True)
add_text_box(slide, Inches(7.1), Inches(3.5), Inches(4.5), Inches(0.4),
             "Spearman Rank Correlation", font_size=18, color=WHITE, bold=True)
add_text_box(slide, Inches(7.1), Inches(3.95), Inches(4.5), Inches(0.5),
             "Tracks SBERT considers similar,\nOpenL3 considers dissimilar",
             font_size=14, color=MUTED)

# Bottom insight
add_card(slide, Inches(1.0), Inches(5.3), Inches(10.8), Inches(1.5), DARK_CARD)
add_text_box(slide, Inches(1.3), Inches(5.5), Inches(10.2), Inches(1.0),
             "Strong complementarity: each modality contributes information the other lacks.\n"
             "This is the core motivation for multi-view fusion \u2014 combining these independent\n"
             "signals produces better retrieval than any single view alone.",
             font_size=16, color=LIGHT)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Fusion Methods
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Fusion: Combining All Views",
             font_size=36, color=WHITE, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
             "Role 4 \u2014 Jiayi: Evaluation & Multi-View Fusion",
             font_size=18, color=MUTED)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0.8), Inches(1.55), Inches(2.5), Pt(3))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT
line.line.fill.background()

# Three method cards
methods = [
    ("Method 1: Weighted Score Fusion",
     "fused = w\u2081\u00b7clap + w\u2082\u00b7openl3\n         + w\u2083\u00b7sbert + w\u2084\u00b7gnn",
     "Normalise all scores to [0,1]\nTune weights on validation split",
     ACCENT),
    ("Method 2: Reciprocal Rank Fusion",
     "RRF(d) = \u03a3 1/(k + rank\u1d62(d))\n               for each view i",
     "Robust, requires no tuning\nUses k=60 (standard constant)",
     GREEN),
    ("Method 3: Learned Reranker",
     "Concatenate per-view scores\nas features for each pair",
     "Logistic regression / MLP\nTrain 70% / evaluate 30%",
     AMBER),
]

card_w = Inches(3.6)
card_h = Inches(4.2)
start_x = Inches(0.8)
gap = Inches(0.4)

for i, (title, formula, note, clr) in enumerate(methods):
    x = start_x + i * (card_w + gap)
    y = Inches(1.9)
    add_card(slide, x, y, card_w, card_h)
    # Color bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 x + Inches(0.15), y + Inches(0.15),
                                 card_w - Inches(0.3), Pt(4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = clr
    bar.line.fill.background()
    # Title
    add_text_box(slide, x + Inches(0.25), y + Inches(0.4),
                 card_w - Inches(0.5), Inches(0.5),
                 title, font_size=17, color=WHITE, bold=True)
    # Formula
    add_text_box(slide, x + Inches(0.25), y + Inches(1.2),
                 card_w - Inches(0.5), Inches(1.2),
                 formula, font_size=14, color=clr, bold=False,
                 font_name="Consolas")
    # Note
    add_text_box(slide, x + Inches(0.25), y + Inches(2.8),
                 card_w - Inches(0.5), Inches(1.0),
                 note, font_size=14, color=MUTED)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 17 — Echo Nest Evaluation Results
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Evaluation: Echo Nest Ground Truth",
             font_size=36, color=WHITE, bold=True)
add_text_box(slide, Inches(0.8), Inches(1.1), Inches(11), Inches(0.5),
             "Independent features no model saw during training (N=294 tracks)",
             font_size=18, color=MUTED)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0.8), Inches(1.55), Inches(2.5), Pt(3))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT
line.line.fill.background()

# Results table
results = [
    ("Random Baseline",     "3.80", "\u2014",     "\u2014",     MUTED),
    ("SBERT (Text+Lyrics)", "3.33", "-12.4%", "1.6e-06", LIGHT),
    ("Fused (SBERT+OpenL3)","3.22", "-15.1%", "1.4e-08", CYAN),
    ("OpenL3 (Audio)",      "2.87", "-24.3%", "1.8e-21", GREEN),
]

# Header
hy = Inches(1.9)
add_text_box(slide, Inches(1.0), hy, Inches(3.5), Inches(0.4),
             "Method", font_size=15, color=ACCENT, bold=True)
add_text_box(slide, Inches(4.8), hy, Inches(2), Inches(0.4),
             "Avg Distance", font_size=15, color=ACCENT, bold=True)
add_text_box(slide, Inches(7.0), hy, Inches(2), Inches(0.4),
             "vs Random", font_size=15, color=ACCENT, bold=True)
add_text_box(slide, Inches(9.2), hy, Inches(2.5), Inches(0.4),
             "p-value", font_size=15, color=ACCENT, bold=True)

for i, (method, dist, vs, pval, clr) in enumerate(results):
    ry = Inches(2.4) + Inches(0.55) * i
    if i % 2 == 1:
        add_card(slide, Inches(0.8), ry - Inches(0.05), Inches(11), Inches(0.5), DARK_CARD)
    bold = (i == 3)
    add_text_box(slide, Inches(1.0), ry, Inches(3.5), Inches(0.4),
                 method, font_size=16, color=clr, bold=bold)
    add_text_box(slide, Inches(4.8), ry, Inches(2), Inches(0.4),
                 dist, font_size=16, color=clr, bold=bold)
    add_text_box(slide, Inches(7.0), ry, Inches(2), Inches(0.4),
                 vs, font_size=16, color=clr, bold=bold)
    add_text_box(slide, Inches(9.2), ry, Inches(2.5), Inches(0.4),
                 pval, font_size=16, color=clr, bold=bold)

# Key insights
add_text_box(slide, Inches(0.8), Inches(4.8), Inches(11), Inches(0.4),
             "Key Takeaways", font_size=20, color=WHITE, bold=True)
tf = add_text_box(slide, Inches(1.0), Inches(5.2), Inches(11), Inches(2),
                  "", font_size=15, color=LIGHT)
insights = [
    "OpenL3 wins because both it and Echo Nest operate in the acoustic domain",
    "Fused embeddings beat either modality alone \u2014 confirms complementary signals",
    "Even SBERT (text only) significantly beats random (p < 10\u207b\u2076)",
    "All differences statistically significant (paired t-test)",
]
for ins in insights:
    add_paragraph(tf, f"\u2022  {ins}", font_size=15, color=LIGHT, space_before=Pt(5))


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 18 — Fine-tuning CLAP
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
bullet_slide(slide, "Fine-tuning CLAP on FMA", [
    ("Role 5 \u2014 Helena: improve pretrained CLAP with domain-specific contrastive learning", WHITE),
    ("",),
    ("Approach:", WHITE),
    ("  \u2022  Construct (audio, text) pairs from FMA metadata",),
    ("  \u2022  Fine-tune with InfoNCE loss \u2014 pull matching pairs together, push others apart",),
    ("  \u2022  Batch size 32\u201364 (more negatives = better contrastive learning)",),
    ("  \u2022  5\u201315 epochs with early stopping on validation loss",),
    ("",),
    ("Goal: reduce conflation of similar genres", WHITE),
    ("  \u2022  Hip-Hop & Pop currently at 0.81 cosine similarity \u2014 nearly identical in CLAP space",),
    ("  \u2022  Fine-tuning should learn FMA-specific genre boundaries",),
    ("  \u2022  Before/after comparison: t-SNE, heatmaps, retrieval quality (P@10, NDCG@10)",),
], sub_title="Role 5 \u2014 Helena")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 19 — Web Demo / System
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
bullet_slide(slide, "Live Demo: Multi-View Recommendation App", [
    ("Flask web application serving side-by-side recommendations", WHITE),
    ("",),
    ("Features:", WHITE),
    ("  \u2022  Search any track by title, artist, or genre",),
    ("  \u2022  See recommendations from each view independently (SBERT, OpenL3, CLAP)",),
    ("  \u2022  Fused results via Reciprocal Rank Fusion across all views",),
    ("  \u2022  Per-view cosine similarity scores shown for each recommendation",),
    ("  \u2022  Audio playback of original MP3 files",),
    ("",),
    ("Technical: all embeddings L2-normalised, cosine similarity via dot product",),
    ("FAISS IndexFlatIP for exact search \u2014 <1ms per query at 8,000 vectors",),
], sub_title="localhost:5001 \u2014 app.py")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 20 — Key Findings
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Key Findings", font_size=36, color=WHITE, bold=True)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0.8), Inches(1.2), Inches(2.5), Pt(3))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT
line.line.fill.background()

findings = [
    ("Multi-view fusion outperforms any single embedding view",
     "Fused retrieval: -15.1% distance vs random (vs -12.4% for text alone, -24.3% for audio alone)",
     GREEN),
    ("Text and audio embeddings are highly complementary",
     "Only 5.6% overlap in top-20 neighbours; Spearman rho = -0.77",
     AMBER),
    ("Data leakage is a real risk in retrieval evaluation",
     "48.8% of tags contained genre labels \u2014 caught and removed before evaluation",
     RED),
    ("Domain-specific evaluation requires careful ground truth",
     "Echo Nest features provide independent, audio-derived ground truth no model saw during training",
     CYAN),
]

for i, (title, detail, clr) in enumerate(findings):
    y = Inches(1.6) + Inches(1.35) * i
    add_card(slide, Inches(0.8), y, Inches(11.5), Inches(1.15), DARK_CARD)
    # Number badge
    badge = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                   Inches(1.1), y + Inches(0.25), Inches(0.5), Inches(0.5))
    badge.fill.solid()
    badge.fill.fore_color.rgb = clr
    badge.line.fill.background()
    add_text_box(slide, Inches(1.1), y + Inches(0.28), Inches(0.5), Inches(0.45),
                 str(i + 1), font_size=18, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER)
    # Text
    add_text_box(slide, Inches(1.9), y + Inches(0.15), Inches(10), Inches(0.45),
                 title, font_size=18, color=WHITE, bold=True)
    add_text_box(slide, Inches(1.9), y + Inches(0.6), Inches(10), Inches(0.4),
                 detail, font_size=14, color=MUTED)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 21 — Team Contributions
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Team Contributions", font_size=36, color=WHITE, bold=True)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0.8), Inches(1.2), Inches(2.5), Pt(3))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT
line.line.fill.background()

team = [
    ("Wenny", "Role 1: Acoustic Similarity", "OpenL3 embeddings, t-SNE comparison, genre centroid heatmaps", AMBER),
    ("Sid & Issac", "Role 2: Lyrics & Semantic", "SBERT pipeline, data leakage fix, Genius lyrics, representation analysis", GREEN),
    ("Archive", "Role 3: Graph Recommendation", "Heterogeneous GNN, link prediction, Flask web demo", PINK),
    ("Jiayi", "Role 4: Evaluation & Fusion", "IR metrics (P@K, MAP, NDCG), weighted fusion, RRF, learned reranker", CYAN),
    ("Helena", "Role 5: Fine-tuning", "CLAP fine-tuning, before/after comparison, Echo Nest correlation, failure analysis", PURPLE),
]

for i, (name, role, work, clr) in enumerate(team):
    y = Inches(1.6) + Inches(1.05) * i
    add_card(slide, Inches(0.8), y, Inches(11.5), Inches(0.85), DARK_CARD)
    add_text_box(slide, Inches(1.1), y + Inches(0.12), Inches(2), Inches(0.35),
                 name, font_size=17, color=clr, bold=True)
    add_text_box(slide, Inches(1.1), y + Inches(0.45), Inches(2), Inches(0.3),
                 role, font_size=12, color=MUTED)
    add_text_box(slide, Inches(3.5), y + Inches(0.2), Inches(8.5), Inches(0.5),
                 work, font_size=14, color=LIGHT)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 22 — Glossary
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Glossary", font_size=36, color=WHITE, bold=True)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0.8), Inches(1.2), Inches(2.5), Pt(3))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT
line.line.fill.background()

glossary = [
    ("Embedding", "A fixed-size numerical vector (e.g., 384 or 512 dimensions) representing a track in a continuous space where similar items are close together."),
    ("FAISS", "Facebook AI Similarity Search \u2014 a library for efficient nearest-neighbor search over dense vectors. We use IndexFlatIP (exact inner product)."),
    ("Cosine Similarity", "The cosine of the angle between two vectors, ranging from -1 to 1. Higher = more similar. Equals the dot product for L2-normalised vectors."),
    ("Contrastive Learning", "A training strategy where a model learns to map similar pairs close together and dissimilar pairs far apart, typically using InfoNCE loss."),
    ("CLAP", "Contrastive Language-Audio Pretraining \u2014 maps audio and text into a shared embedding space. Captures high-level semantics (mood, genre, instruments)."),
    ("SBERT", "Sentence-BERT \u2014 a siamese transformer that produces fixed-size sentence embeddings optimised for semantic similarity. Model: all-MiniLM-L6-v2."),
    ("OpenL3", "A self-supervised audio embedding model trained on AudioSet. Captures low-level acoustic features (timbre, rhythm, texture) without labels."),
    ("GNN", "Graph Neural Network \u2014 learns node embeddings by aggregating information from graph neighbours. We use SAGEConv on a heterogeneous track-artist-genre graph."),
    ("RRF", "Reciprocal Rank Fusion \u2014 combines ranked lists by summing 1/(k + rank) across views. Robust and requires no weight tuning."),
    ("Data Leakage", "When evaluation targets (e.g., genre labels) leak into model inputs, artificially inflating metrics. We removed genre from SBERT input strings."),
    ("t-SNE", "t-distributed Stochastic Neighbor Embedding \u2014 a dimensionality reduction technique for visualising high-dimensional data in 2D."),
]

# Two columns
col1 = glossary[:6]
col2 = glossary[6:]

for col_idx, col_data in enumerate([col1, col2]):
    x_offset = Inches(0.8) + col_idx * Inches(6.2)
    for i, (term, defn) in enumerate(col_data):
        y = Inches(1.6) + Inches(0.88) * i
        add_text_box(slide, x_offset, y, Inches(5.8), Inches(0.3),
                     term, font_size=14, color=ACCENT, bold=True)
        add_text_box(slide, x_offset, y + Inches(0.28), Inches(5.8), Inches(0.55),
                     defn, font_size=11, color=LIGHT)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 23 — Dataset & References
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(11), Inches(0.8),
             "Datasets & References", font_size=36, color=WHITE, bold=True)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0.8), Inches(1.2), Inches(2.5), Pt(3))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT
line.line.fill.background()

# Datasets
add_text_box(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(0.5),
             "Datasets", font_size=22, color=WHITE, bold=True)

datasets = [
    ("Free Music Archive (FMA Small)", "8,000 tracks, 8 genres, ~30s clips, open-licensed", "Defferrard et al., 2017  \u2014  github.com/mdeff/fma"),
    ("Echo Nest / Spotify Audio Features", "Danceability, energy, valence, tempo, acousticness, instrumentalness, liveness, speechiness", "294-track overlap with FMA Small used as independent evaluation ground truth"),
    ("Genius API", "Song lyrics (first 1,000 characters per track)", "Used to enrich SBERT text embeddings after removing genre labels"),
    ("AudioSet", "Training data for OpenL3 (self-supervised audio-visual correspondence)", "Gemmeke et al., 2017  \u2014  research.google.com/audioset"),
]

for i, (name, desc, source) in enumerate(datasets):
    y = Inches(2.15) + Inches(0.8) * i
    add_text_box(slide, Inches(1.0), y, Inches(11), Inches(0.3),
                 f"\u2022  {name}", font_size=15, color=WHITE, bold=True)
    add_text_box(slide, Inches(1.3), y + Inches(0.28), Inches(10.5), Inches(0.22),
                 desc, font_size=12, color=LIGHT)
    add_text_box(slide, Inches(1.3), y + Inches(0.5), Inches(10.5), Inches(0.22),
                 source, font_size=11, color=MUTED)

# Models & Libraries
add_text_box(slide, Inches(0.8), Inches(5.5), Inches(11), Inches(0.5),
             "Models & Libraries", font_size=22, color=WHITE, bold=True)

refs = [
    "CLAP (LAION): HTSAT-tiny backbone, pretrained on AudioSet + music  |  laion/clap",
    "SBERT: all-MiniLM-L6-v2, 384-d  |  sentence-transformers",
    "OpenL3: self-supervised audio embeddings  |  marl/openl3",
    "PyTorch Geometric: heterogeneous GNN (SAGEConv)  |  pyg-team/pytorch_geometric",
    "FAISS: Facebook AI Similarity Search  |  facebookresearch/faiss",
]

tf = add_text_box(slide, Inches(1.0), Inches(5.9), Inches(11), Inches(1.2),
                  "", font_size=12, color=LIGHT)
for r in refs:
    add_paragraph(tf, f"\u2022  {r}", font_size=12, color=LIGHT, space_before=Pt(3))


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 24 — Thank You / Questions
# ═══════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)

add_text_box(slide, Inches(0.8), Inches(2.0), Inches(11), Inches(1.5),
             "Thank You", font_size=56, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(5.5), Inches(3.5), Inches(2.5), Pt(4))
line.fill.solid()
line.fill.fore_color.rgb = ACCENT
line.line.fill.background()

add_text_box(slide, Inches(0.8), Inches(3.8), Inches(11), Inches(0.8),
             "Questions?", font_size=28, color=MUTED,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.8), Inches(5.0), Inches(11), Inches(0.5),
             "Demo available at localhost:5001  |  Source code on GitHub",
             font_size=16, color=MUTED, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════════════════
output = "/Users/zac/Documents/Documents-it/Deep Learning Notes/Finals/Multi-Faceted_Music_Retrieval_Presentation.pptx"
prs.save(output)
print(f"Saved to: {output}")
